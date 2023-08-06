import sys
import os
import matplotlib.pyplot as plt
from pptx import Presentation
import io 
from PyQt5 import QtWidgets
from PyQt5.uic import loadUiType 
from matplotlib.figure import Figure
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.backends.backend_qt5agg import (
    FigureCanvasQTAgg as FigureCanvas,
    NavigationToolbar2QT as NavigationToolbar)

#get absolute path of QT5 UI
path = os.path.abspath(os.path.dirname(__file__))
#import parent classes from QT5 UI
Ui_MainWindow, QMainWindow = loadUiType(f"{path}/data/window.ui")

class Main(QMainWindow, Ui_MainWindow):
    
    def __init__(self, ):
        super(Main, self).__init__()
        self.setupUi(self)
        self.fig_dict = {}
        self.aliases_dict = {}
        self.titles = []
        self.axes = []
        self.fignums = []
        #change fig when list widget selection changes
        self.listWidget.currentRowChanged.connect(self.change_fig)
        #populate options for comboBox figure name aliases
        for option in ["Title","Axes","Fig. Num"]:
            self.comboBox.addItem(option)
        #change fig name aliases when combobox changes    
        self.comboBox.currentTextChanged.connect(self.change_figname)
        #save plots to file on button click
        self.saveButton.clicked.connect(self.file_save)
        #inject empty figure at start
        fig = Figure()
        self.add_canvas(fig)
 
    def add_canvas(self, fig): 
        """
        Add desired Plot and toolbar canvas from the generic window widget
        vertical layout (canvas_vlayout) when plot is changed
        """
        self.canvas = FigureCanvas(fig)
        self.toolbar = NavigationToolbar(self.canvas, 
                self, coordinates=True)
        self.canvas_vlayout.addWidget(self.toolbar)
        self.canvas_vlayout.addWidget(self.canvas)
        self.canvas.draw()
               
    def remove_canvas(self,):
        """
        Remove current Plot and toolbar canvas from the generic window widget
        vertical layout (canvas_vlayout) when plot is changed
        """
        #reset plot view beofre change
        self.canvas.toolbar.home()
        #remove widgets from canvas_vlayout
        self.canvas_vlayout.removeWidget(self.toolbar)
        self.toolbar.close() 
        self.canvas_vlayout.removeWidget(self.canvas)
        self.canvas.close()
              
    def append_fig(self, title, axis, fignum, name, fig):
        """
        Build dictionaries that will deliver the desired figure based
        on selected alias 
        """
        #main dictionary
        self.fig_dict[name] = fig
        #aliases dictionaries
        self.aliases_dict[title] = name
        self.aliases_dict[axis] = name
        self.aliases_dict[fignum] = name
        #figure list aliases entries
        self.titles.append(title)
        self.axes.append(axis)
        self.fignums.append(fignum)
        #figure list default entries
        self.listWidget.addItem(title)  
        
    def change_fig(self, item):
        """
        inject selected figure to canvas
        """
        text = self.listWidget.currentItem().text()
        self.remove_canvas()
        self.add_canvas(self.fig_dict[self.aliases_dict[text]]) 
        
    def change_figname(self, selection):
        """
        change the plot identifier text in listWidget
        """
        #get desired aliases list
        if selection == "Title":
            options = self.titles
        elif selection == "Axes":
            options = self.axes
        else:
            options = self.fignums
        #replace current list with desired aliases 
        for index, option in enumerate(options):
            item = self.listWidget.item(index) 
            item.setText(option)
            
    def file_save(self):
        """
        saves collected matplotlib figures in a .pptx or pdf file
        """
        #collect figure list to be saved to file
        figures = list(map(plt.figure,plt.get_fignums()))
        #collect save_path and desired file extension
        extensions = "PowerPoint Presentation(*.pptx);;PDF(*.pdf)"
        save_path, extension = QtWidgets.QFileDialog.getSaveFileName(self, 'Save File',"Figures",extensions)  
        # save path in desired format
        if extension == "PDF(*.pdf)":
            with PdfPages(save_path) as pdf:
                for fig in figures:
                    pdf.savefig(fig,bbox_inches='tight',pad_inches=0.1)#save pdf
            QtWidgets.QMessageBox.about(self, "Save File", "PDF file saved.")
            
        elif extension == "PowerPoint Presentation(*.pptx)":
            #initialize presentation object
            prs = Presentation()
            #force slide aspect ratio to widescreen
            prs.slide_width = 11887200
            prs.slide_height = 6686550
            
            for fig in figures:
                #save figure to file like object
                image_stream = io.BytesIO()
                fig.savefig(image_stream,bbox_inches='tight',pad_inches=0.1)
                #build blank slide and insert figure into
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                left = top = 1
                slide.shapes.add_picture(image_stream,left,top)
                #close figure stream
                image_stream.close()

            prs.save(save_path) #save pptx
            QtWidgets.QMessageBox.about(self, "Save File", "PowerPoint Presentation saved.")
        else:
            pass
        
def view():
    """
    Collects all matplotlib figures and displays them in a Qt5 GUI
    """
    #collect figures in list
    figures = list(map(plt.figure,plt.get_fignums()))
    #start app
    app = QtWidgets.QApplication(sys.argv)
    main = Main()
    
    if figures:
        for count, figure in enumerate(figures):
            #main names for figures
            name   = f"{figure.number}"
            #aliases for figures
            titles  = [figure.axes[0].get_title(loc = i) for i in ["left","center","right"]]
            titles = [i for i in titles if i]
            title = f"{count+1}- {titles[0]}" if titles else ""
            axes_labels   = f"{count+1}- {figure.axes[0].get_ylabel()} vs {figure.axes[0].get_xlabel()} "
            fignum = f"Figure {figure.number}"
            #Append figure to App
            main.append_fig(title, axes_labels, fignum, name, figure)
            
    main.show()
    sys.exit(app.exec_())
 
if __name__ == '__main__':
    import numpy as np
    plt.ioff()
    
    x = np.linspace(0, 2*np.pi, 30)
    for c, i in enumerate(range(1, 5)):
        y = np.sin(x+i)
        plt.figure(f"figure Sin omega + {c}",figsize = (13,6.5),dpi=100)
        plt.plot(x, y, marker="o", label=r"$sin(\theta) + \omega$"+f"$_{c}$")
        plt.xlabel(r"$\theta$")
        plt.ylabel(f"Sin of theta + omega {c}")
        plt.xlim(0, 2*np.pi)
        plt.title(f"Sin of theta + omega {c} \nexample", loc="center")
        plt.legend(bbox_to_anchor=(1.01, 1.0),
                         loc='upper left', title="Functions")
        
        plt.tight_layout()
        plt.grid()
    view()
    
     
    